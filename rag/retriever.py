import logging
import sys
from pythonjsonlogger import jsonlogger

import os, uuid, shutil
from typing import List, Dict, Any
import pypdf
import docx
import pptx

from fastapi import FastAPI, UploadFile, File, HTTPException, Depends, Header
from pydantic import BaseModel
from qdrant_client import QdrantClient, models
from qdrant_client.http.models import VectorParams, Distance, PointStruct, UpdateStatus
from sentence_transformers import SentenceTransformer

# OpenTelemetry Imports
from opentelemetry import trace
from opentelemetry.sdk.trace import TracerProvider, NoOpTracer
from opentelemetry.sdk.trace.export import BatchSpanProcessor
from opentelemetry.exporter.otlp.proto.grpc.trace_exporter import OTLPSpanExporter
from opentelemetry.instrumentation.fastapi import FastAPIInstrumentor

# --- Logging Setup ---
# Configure root logger for structured JSON logging
logHandler = logging.StreamHandler(sys.stdout)
formatter = jsonlogger.JsonFormatter(
    fmt="%(asctime)s %(name)s %(levelname)s %(message)s"
)
logHandler.setFormatter(formatter)
logging.basicConfig(handlers=[logHandler], level=logging.INFO)
logging.getLogger("uvicorn.access").disabled = True  # Disable default access logs
# --- End Logging Setup ---

# --- Tracing Setup ---
API_KEY = os.getenv("BEND_API_KEY")
OTEL_ENDPOINT = os.getenv("OTEL_EXPORTER_OTLP_ENDPOINT")
tracer: trace.Tracer

if OTEL_ENDPOINT:
    provider = TracerProvider()
    processor = BatchSpanProcessor(
        OTLPSpanExporter(endpoint=OTEL_ENDPOINT, insecure=True)
    )
    provider.add_span_processor(processor)
    trace.set_tracer_provider(provider)
    tracer = trace.get_tracer(__name__)
else:
    # Use a No-op tracer if the endpoint isn't configured
    tracer = NoOpTracer()
# --- End Tracing Setup ---


async def api_key_security(x_api_key: str = Header(None)):
    if API_KEY:  # Security is enabled
        if x_api_key != API_KEY:
            raise HTTPException(status_code=403, detail="Invalid API Key")


app = FastAPI(dependencies=[Depends(api_key_security)])

if OTEL_ENDPOINT:
    FastAPIInstrumentor.instrument_app(app)

client = QdrantClient(host="qdrant", port=6333)
model = SentenceTransformer("all-MiniLM-L6-v2")
COLLECTION = "bend-docs"

try:
    client.get_collection(collection_name=COLLECTION)
except Exception:
    client.create_collection(
        collection_name=COLLECTION,
        vectors_config=VectorParams(
            size=model.get_sentence_embedding_dimension(),
            distance=Distance.COSINE,
        ),
    )


def _split_text(text: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    """A simple, dependency-free text splitter."""
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - chunk_overlap
        if start >= len(text):
            break
    return chunks


def _embed_and_store_text(text: str, source: str):
    with tracer.start_as_current_span("embed_and_store_text"):
        with tracer.start_as_current_span("chunking_text"):
            chunks = _split_text(text, chunk_size=500, chunk_overlap=50)

        with tracer.start_as_current_span("encoding_vectors"):
            vectors = model.encode(chunks).tolist()

        points = [
            PointStruct(
                id=str(uuid.uuid4()),
                vector=vec,
                payload={"text": chunk, "source": source},
            )
            for vec, chunk in zip(vectors, chunks)
        ]

        with tracer.start_as_current_span("upsert_to_qdrant"):
            client.upsert(collection_name=COLLECTION, points=points, wait=True)


@app.post("/ingest")
async def ingest_doc(file: UploadFile = File(...)):
    try:
        filename = file.filename
        text = ""

        if filename.lower().endswith(".pdf"):
            reader = pypdf.PdfReader(file.file)
            for page in reader.pages:
                text += page.extract_text()
        elif filename.lower().endswith(".docx"):
            document = docx.Document(file.file)
            for para in document.paragraphs:
                text += para.text + "\n"
        elif filename.lower().endswith(".pptx"):
            presentation = pptx.Presentation(file.file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif filename.lower().endswith((".txt", ".md")):
            file_content = await file.read()
            text = file_content.decode("utf-8")
        else:
            raise HTTPException(
                status_code=400,
                detail="Unsupported file type. Please upload a .pdf, .docx, .pptx, .txt, or .md file.",
            )

        _embed_and_store_text(text, filename)
        return {"message": f"'{filename}' ingested successfully"}
    except Exception as e:
        # Catch potential parsing errors and other issues
        raise HTTPException(status_code=500, detail=f"Failed to ingest document: {e}")


class IngestTextRequest(BaseModel):
    text: str
    source: str


@app.post("/ingest/text")
async def ingest_text(req: IngestTextRequest):
    try:
        _embed_and_store_text(req.text, req.source)
        return {"message": f"Text from source '{req.source}' ingested successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to ingest text: {e}")


class QueryRequest(BaseModel):
    query: str
    top_k: int = 5


@app.post("/retrieve")
def retrieve(req: QueryRequest):
    try:
        with tracer.start_as_current_span("retrieve_from_qdrant"):
            vec = model.encode(req.query).tolist()
            hits = client.search(
                collection_name=COLLECTION, query_vector=vec, limit=req.top_k
            )
            return [hit.payload for hit in hits]
    except Exception as e:
        raise HTTPException(
            status_code=503, detail=f"Failed to retrieve from vector store: {e}"
        )


@app.get("/documents")
def get_documents():
    try:
        sources = set()
        response = client.scroll(
            collection_name=COLLECTION,
            with_payload=["source"],
            with_vectors=False,
            limit=1000,
        )
        for point in response[0]:
            if point.payload and "source" in point.payload:
                sources.add(point.payload["source"])
        return list(sources)
    except Exception as e:
        raise HTTPException(status_code=503, detail=f"Failed to list documents: {e}")


class DeleteDocRequest(BaseModel):
    source: str


@app.delete("/documents")
def delete_document(req: DeleteDocRequest):
    try:
        with tracer.start_as_current_span("delete_from_qdrant"):
            result = client.delete(
                collection_name=COLLECTION,
                points_selector=models.FilterSelector(
                    filter=models.Filter(
                        must=[
                            models.FieldCondition(
                                key="source",
                                match=models.MatchValue(value=req.source),
                            )
                        ]
                    )
                ),
                wait=True,
            )
        if result.status == UpdateStatus.COMPLETED:
            return {
                "message": f"Document with source '{req.source}' deleted successfully."
            }
        else:
            raise HTTPException(
                status_code=500,
                detail=f"Failed to delete document. Status: {result.status}",
            )
    except Exception as e:
        raise HTTPException(
            status_code=503, detail=f"Failed to delete document from vector store: {e}"
        )


class OpenWebUIRequest(BaseModel):
    query: str
    chat_id: str
    headers: Dict[str, Any]


@app.post("/openwebui_rag", dependencies=[Depends(api_key_security)])
def openwebui_rag(req: OpenWebUIRequest):
    try:
        with tracer.start_as_current_span("openwebui_rag_retrieval"):
            vec = model.encode(req.query).tolist()
            hits = client.search(collection_name=COLLECTION, query_vector=vec, limit=3)

        if not hits:
            return {"content": ""}

        context_str = "CONTEXT:\n"
        for hit in hits:
            context_str += (
                f"- Source: {hit.payload['source']}\n  Content: {hit.payload['text']}\n"
            )

        return {"content": context_str}
    except Exception as e:
        logging.error(f"Error during OpenWebUI RAG retrieval: {e}", exc_info=True)
        return {"content": ""}
